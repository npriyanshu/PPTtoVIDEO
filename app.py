from flask import Flask, request, jsonify, render_template, url_for, send_from_directory
from pptx import Presentation
from gtts import gTTS
from flask_cors import CORS
from moviepy.editor import concatenate_videoclips, ImageClip, AudioFileClip
from PIL import Image, ImageDraw, ImageFont
from io import BytesIO
import numpy as np
import os

app = Flask(__name__)
CORS(app, origins='*', allow_headers=['Content-Type'])

@app.route('/')
def index():
    return render_template('index.html')

@app.route('/convert', methods=['POST'])
def convert_ppt_to_video():
    print("Received a request to convert PPT to video")
    file = request.files['file']
    if not file:
        return jsonify({"error": "No file provided"}), 400

    # Create a directory to store intermediate files
    media_dir = os.path.abspath('media')
    if not os.path.exists(media_dir):
        os.makedirs(media_dir)

    # Load the PowerPoint presentation
    prs = Presentation(file)

    # List to hold the video clips
    clips = []
    audio_paths = []

    # Generate video clips from each slide
    for i, slide in enumerate(prs.slides):
        # Convert slide to image
        slide_image_path = os.path.join(media_dir, f"slide_{i}.png")
        slide_image = slide_to_image(slide)
        slide_image.save(slide_image_path)
        
        # Convert the image to a NumPy array
        img = Image.open(slide_image_path)
        img_np = np.array(img)

        # Create an ImageClip from the slide image
        img_clip = ImageClip(img_np)

        # Extract text from the slide for audio
        text = extract_text_from_slide(slide)
        if not text:
            text = "No text found on this slide."

        # Generate audio for the slide using gTTS
        audio_path = os.path.join(media_dir, f"audio_{i}.mp3")
        tts = gTTS(text=text, lang='en')
        tts.save(audio_path)
        audio_clip = AudioFileClip(audio_path)
        audio_paths.append(audio_path)

        # Set the duration of the image clip to match the audio clip
        img_clip = img_clip.set_duration(audio_clip.duration)

        # Set the audio of the image clip
        img_clip = img_clip.set_audio(audio_clip)

        # Add the clip to the list
        clips.append(img_clip)

    # Concatenate all the clips into one video
    video = concatenate_videoclips(clips, method="compose")

    # Ensure the static directory exists
    os.makedirs('static', exist_ok=True)

    # Save the video to a file
    video_path = "static/output_video.mp4"
    video.write_videofile(video_path, fps=24)

    # Cleanup temporary audio files
    for path in audio_paths:
        os.remove(path)

    # Cleanup slide image files
    for i in range(len(prs.slides)):
        os.remove(os.path.join(media_dir, f"slide_{i}.png"))

    return jsonify({"video_url": url_for('static', filename='output_video.mp4')})

def slide_to_image(slide):
    # Create a blank image with a white background
    img = Image.new('RGB', (1920, 1080), color='white')
    draw = ImageDraw.Draw(img)
    font = ImageFont.truetype("arial.ttf", 24)  # Use a truetype font

    # Draw each shape's text onto the image
    for shape in slide.shapes:
        if shape.has_text_frame:
            for paragraph in shape.text_frame.paragraphs:
                for run in paragraph.runs:
                    # Calculate position (you might need to adjust x, y)
                    x, y = shape.left // 12700, shape.top // 12700  # Convert EMU to pixels
                    draw.text((x, y), run.text, fill="black", font=font)

        elif shape.shape_type == 13:  # Shape type 13 corresponds to pictures
            image_stream = shape.image.blob
            image = Image.open(BytesIO(image_stream))
            image = image.resize((shape.width // 12700, shape.height // 12700), Image.LANCZOS)  # Convert EMU to pixels
            img.paste(image, (shape.left // 12700, shape.top // 12700))

    return img

def extract_text_from_slide(slide):
    text = ""
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        for paragraph in shape.text_frame.paragraphs:
            for run in paragraph.runs:
                text += run.text + " "
    return text.strip()

if __name__ == "__main__":
    app.run(debug=True)




# from flask import Flask, request, jsonify, render_template, url_for, send_from_directory
# from pptx import Presentation
# from gtts import gTTS
# from flask_cors import CORS
# from moviepy.editor import concatenate_videoclips, ImageClip, AudioFileClip
# from PIL import Image, ImageDraw, ImageFont
# from io import BytesIO
# import numpy as np
# import os

# app = Flask(__name__)
# CORS(app, origins='*', allow_headers=['Content-Type'])

# @app.route('/')
# def index():
#     return render_template('index.html')

# @app.route('/convert', methods=['POST'])
# def convert_ppt_to_video():
#     print("Received a request to convert PPT to video")
#     file = request.files['file']
#     if not file:
#         return jsonify({"error": "No file provided"}), 400

#     # Load the PowerPoint presentation
#     prs = Presentation(file)

#     # List to hold the video clips
#     clips = []
#     audio_paths = []

#     # Generate video clips from each slide
#     for i, slide in enumerate(prs.slides):
#         # Convert slide to image
#         img_stream = BytesIO()
#         slide_image = slide_to_image(slide)
#         slide_image.save(img_stream, format='PNG')
#         img_stream.seek(0)

#         # Convert the image to a NumPy array
#         img = Image.open(img_stream)
#         img_np = np.array(img)

#         # Create an ImageClip from the slide image
#         img_clip = ImageClip(img_np)

#         # Extract text from the slide for audio
#         text = extract_text_from_slide(slide)
#         if not text:
#             text = "No text found on this slide."

#         # Generate audio for the slide using gTTS
#         tts = gTTS(text=text, lang='en')
#         audio_path = f"audio_{i}.mp3"
#         tts.save(audio_path)
#         audio_clip = AudioFileClip(audio_path)
#         audio_paths.append(audio_path)

#         # Set the duration of the image clip to match the audio clip
#         img_clip = img_clip.set_duration(audio_clip.duration)

#         # Set the audio of the image clip
#         img_clip = img_clip.set_audio(audio_clip)

#         # Add the clip to the list
#         clips.append(img_clip)

#     # Concatenate all the clips into one video
#     video = concatenate_videoclips(clips, method="compose")

#     # Ensure the static directory exists
#     os.makedirs('static', exist_ok=True)

#     # Save the video to a file
#     video_path = "static/output_video.mp4"
#     video.write_videofile(video_path, fps=24)

#     # Cleanup temporary audio files
#     for path in audio_paths:
#         os.remove(path)

#     return jsonify({"video_url": url_for('static', filename='output_video.mp4')})

# def slide_to_image(slide):
#     # Create a blank image with a white background
#     img = Image.new('RGB', (1920, 1080), color='white')
#     draw = ImageDraw.Draw(img)
#     font = ImageFont.truetype("arial.ttf", 24)  # Use a truetype font

#     # Draw each shape's text onto the image
#     for shape in slide.shapes:
#         if shape.has_text_frame:
#             for paragraph in shape.text_frame.paragraphs:
#                 for run in paragraph.runs:
#                     # Calculate position (you might need to adjust x, y)
#                     x, y = shape.left // 12700, shape.top // 12700  # Convert EMU to pixels
#                     draw.text((x, y), run.text, fill="black", font=font)

#         elif shape.shape_type == 13:  # Shape type 13 corresponds to pictures
#             image_stream = shape.image.blob
#             image = Image.open(BytesIO(image_stream))
#             image = image.resize((shape.width // 12700, shape.height // 12700), Image.LANCZOS)  # Convert EMU to pixels
#             img.paste(image, (shape.left // 12700, shape.top // 12700))

#     return img


# def extract_text_from_slide(slide):
#     text = ""
#     for shape in slide.shapes:
#         if not shape.has_text_frame:
#             continue
#         for paragraph in shape.text_frame.paragraphs:
#             for run in paragraph.runs:
#                 text += run.text + " "
#     return text.strip()

# if __name__ == "__main__":
#     app.run(debug=True)

